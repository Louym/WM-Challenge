#!/usr/bin/env python3
"""Generate slides.pptx from the manifests, embedding compressed videos from pptx_assets/."""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).parent
ASSETS = ROOT / "pptx_assets"
SAMPLES = ["demo_29", "demo_42", "demo_79"]
SAMPLE_LABELS = {"demo_29": "sample 29 (coffee shop)", "demo_42": "sample 42 (car dashboard)", "demo_79": "sample 79"}

BG = RGBColor(0x0E, 0x11, 0x16)
FG = RGBColor(0xE8, 0xED, 0xF2)
ACCENT = RGBColor(0x7F, 0xC4, 0xFF)
DIM = RGBColor(0x93, 0xA1, 0xB0)

MODEL_BLURB = {
    "HY-WorldPlay": "Tencent HY-WorldPlay on HunyuanVideo-1.5 (480p i2v, 832x480 @ 24fps). Four checkpoints: ar_distill (4-step distilled causal), ar (20-step causal), ar_rl (RL-tuned causal), bi (bidirectional). Camera/action conditioning via pose trajectories.",
    "Sana-WM": "Sana-WM bidirectional 1600M 720p (1280x704 @ 16fps) with refiner; 60 sampling steps, CFG 5, Pi3X-estimated intrinsics. Actions are camera trajectories (c2w) from a WASD/IJKL DSL.",
    "LingBot-World": "LingBot-World i2v-A14B (832x480 @ 16fps). small/fast = distilled chunked-streaming causal pipeline; large/base = 20-step act2cam pipeline. Dual (high/low-noise) 14B Wan-style models.",
    "LingBot-World-v2": "LingBot-World 2.0 / LingBot-World-Infinity (robbyant, arXiv:2607.07534), built on Wan2.2 i2v-A14B (832x480 @ 16fps). Causal pretraining for an unbounded interaction horizon; released variant is the 4-step distilled causal-fast model streaming with a KV window (local_attn 18, sink 6). Code: github.com/robbyant/lingbot-world-v2.",
    "Matrix-Game-3.0": "Skywork Matrix-Game-3.0 interactive streaming world model (1280x704 @ 17fps). efficiency = distilled+int8, quality = base model at 50 steps. One action per 40-frame chunk, KV state persists across the stream (runs are one continuous causal video; no whole-video mode). Run 2-GPU distributed.",
    "DreamX-World-5B": "DreamX-World-5B autoregressive long-horizon variant on Wan2.2-TI2V-5B (1280x704 @ 16fps). One camera+move token per generation; 21 latent frames per clip.",
    "DreamX-World-5B-Cam": "DreamX-World-5B-Cam bidirectional camera-controlled variant (1280x704 @ 16fps), 50 steps, CFG 3. Whole-video run is ~12x its training clip length and needed CPU offload.",
    "Infinite-World": "MeiGen Infinite-World chunk-autoregressive model (~627p bucket @ 16fps), 30 steps, CFG 5. Fixed 49-frame chunks per action; no whole-video mode.",
}


def load_manifests():
    per_sample = {}
    for sample in SAMPLES:
        suf = "" if sample == "demo_29" else f"_{sample}"
        per_sample[sample] = {m["video"].replace(".mp4", ""): m for m in json.loads((ROOT / f"manifest{suf}.json").read_text())}
    return per_sample


def dark_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_text(slide, x, y, w, h, text, size, color=FG, bold=False, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = line
        f = r.font
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
    return box


def main() -> None:
    manifests = load_manifests()
    base = manifests["demo_29"]
    runs = list(base.keys())

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # title
    s = dark_slide(prs)
    add_text(s, 0.8, 2.1, 11.7, 1.2, "World-Model Baselines — Interactive Demo Comparison", 38, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 3.4, 11.7, 0.6, "7 models x all supported settings (causal vs bidirectional, clip-by-clip vs whole-video)", 20, align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 4.0, 11.7, 0.6, "3 input samples (29, 42, 79)  ·  69 videos  ·  ~60 s each  ·  shared trajectory", 20, align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 4.9, 11.7, 0.5, "Generated with the unmodified baselines/demo server driven over its HTTP API  ·  H200 GPUs", 14, color=DIM, align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 5.4, 11.7, 0.5, "Videos are embedded at 480p (click to play). Full-resolution files: WM-Challenge repo.", 14, color=DIM, align=PP_ALIGN.CENTER)

    # protocol
    s = dark_slide(prs)
    add_text(s, 0.7, 0.5, 12, 0.8, "Protocol", 30, color=ACCENT, bold=True)
    add_text(s, 0.9, 1.5, 11.5, 5.5,
             "Inputs: SANA_WM_dev/playground/dataset samples demo_29 (coffee shop), demo_42 (car dashboard), demo_79 — image + text prompt.\n"
             "\n"
             "Shared camera trajectory (~60 s): forward → pan right → forward → backward, equal quarters. "
             "DreamX whole-video runs are push-in only (its API takes one action token per shot).\n"
             "\n"
             "clip-by-clip: one action per segment; each segment restarts from the last frame of the previous one, so temporal context is lost at boundaries. "
             "Exceptions: Matrix-Game-3 and Infinite-World stream continuously with persistent state.\n"
             "\n"
             "whole-video: a single action string / one long rollout in one generation call.\n"
             "\n"
             "Latency shown per video is generation wall-time on H200 GPUs, excluding model load. Full numbers in latency*.csv.",
             16)

    current_model = None
    for run in runs:
        cfg = base[run]
        model = cfg["model"]
        if model != current_model:
            current_model = model
            s = dark_slide(prs)
            add_text(s, 0.7, 0.6, 12, 0.9, model, 34, color=ACCENT, bold=True)
            add_text(s, 0.9, 1.8, 11.5, 2.2, MODEL_BLURB.get(model, ""), 17)
            variants = sorted({base[r]["variant"] for r in runs if base[r]["model"] == model})
            add_text(s, 0.9, 4.3, 11.5, 2.2,
                     f"variants run: {', '.join(variants)}\n"
                     f"temporal: {cfg['temporal']}\n"
                     f"resolution / fps: {cfg.get('measured_resolution', cfg['resolution'])} @ {cfg['fps']} fps",
                     16, color=DIM)

        s = dark_slide(prs)
        add_text(s, 0.5, 0.35, 12.4, 0.6, f"{model}  ·  {cfg['variant']}  —  {cfg['protocol']}", 20, bold=True)

        # 3 videos side by side, 16:9-ish boxes
        vw, vh = 4.05, 2.55
        gap = 0.15
        x0 = (13.333 - 3 * vw - 2 * gap) / 2
        y0 = 1.15
        for i, sample in enumerate(SAMPLES):
            vid = ASSETS / f"{sample}__{run}.mp4"
            poster = ASSETS / f"{sample}__{run}.jpg"
            x = x0 + i * (vw + gap)
            if vid.exists():
                s.shapes.add_movie(str(vid), Inches(x), Inches(y0), Inches(vw), Inches(vh),
                                   poster_frame_image=str(poster) if poster.exists() else None,
                                   mime_type="video/mp4")
            m = manifests[sample].get(run, {})
            gen = m.get("total_generation_s")
            gen_str = f"{gen/60:.1f} min gen" if gen else "n/a"
            add_text(s, x, y0 + vh + 0.08, vw, 0.4,
                     f"{SAMPLE_LABELS[sample]}  ·  {m.get('duration_s', '?')} s  ·  {gen_str}",
                     11, color=DIM, align=PP_ALIGN.CENTER)

        steps = f"  ·  {cfg['steps']} steps" if cfg.get("steps") else ""
        cfgs = f"  ·  CFG {cfg['cfg']}" if cfg.get("cfg") else ""
        add_text(s, 0.7, 4.75, 12, 0.9,
                 f"{cfg['temporal']}  ·  {cfg['clips']} x {cfg['frames_per_clip']} frames @ {cfg['fps']} fps{steps}{cfgs}  ·  seed {cfg['seed']}",
                 13)
        add_text(s, 0.7, 5.45, 12, 0.6, cfg["trajectory"], 12, color=DIM)
        add_text(s, 0.7, 6.0, 12, 0.9, cfg["notes"], 12, color=DIM, italic=True)

    # latency table
    s = dark_slide(prs)
    add_text(s, 0.7, 0.35, 12, 0.7, "Generation latency (minutes, H200, excl. model load)", 24, color=ACCENT, bold=True)
    n_rows = len(runs) + 1
    tbl = s.shapes.add_table(n_rows, 4, Inches(1.6), Inches(1.15), Inches(10.1), Inches(0.28 * n_rows)).table
    for j, h in enumerate(["run", "sample 29", "sample 42", "sample 79"]):
        c = tbl.cell(0, j)
        c.text = h
    for i, run in enumerate(runs, start=1):
        tbl.cell(i, 0).text = run
        for j, sample in enumerate(SAMPLES, start=1):
            g = manifests[sample].get(run, {}).get("total_generation_s")
            tbl.cell(i, j).text = f"{g/60:.1f}" if g else "-"
    for row in tbl.rows:
        row.height = Emu(int(Inches(0.27)))
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.color.rgb = FG
    add_text(s, 0.7, 7.0, 12, 0.4,
             "Multi-GPU runs: mg3_* (2 GPUs), lingbot_large_whole (4 GPUs), sample-29 dreamx_cam_whole (CPU offload).",
             11, color=DIM)

    # Autoplay: PowerPoint defaults embedded media to click-to-play via
    # <p:cond delay="indefinite"/> in each video's timing node; delay="0"
    # starts playback when the slide appears.
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    for slide in prs.slides:
        for cond in slide._element.findall(".//p:video//p:cond", ns):
            if cond.get("delay") == "indefinite":
                cond.set("delay", "0")

    out = ROOT / "slides.pptx"
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides.slides if hasattr(prs.slides, 'slides') else prs.slides._sldIdLst)} slides, {out.stat().st_size/1e6:.0f} MB)")


if __name__ == "__main__":
    main()
